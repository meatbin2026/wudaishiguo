from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


TITLE_FONT = "Microsoft YaHei"
BODY_FONT = "Microsoft YaHei"

NAVY = RGBColor(15, 39, 74)
TEAL = RGBColor(9, 122, 122)
MINT = RGBColor(215, 241, 238)
LIGHT_BG = RGBColor(245, 248, 250)
MID_GRAY = RGBColor(88, 102, 119)
DARK = RGBColor(33, 39, 46)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(210, 220, 230)
GREEN = RGBColor(45, 168, 118)
ORANGE = RGBColor(245, 146, 69)
RED = RGBColor(226, 86, 72)


def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def add_footer(slide, page, total):
    box = slide.shapes.add_textbox(Inches(11.5), Inches(7.08), Inches(1.6), Inches(0.25))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{page:02d}/{total:02d}"
    p.alignment = PP_ALIGN.RIGHT
    run = p.runs[0]
    run.font.name = BODY_FONT
    run.font.size = Pt(11)
    run.font.color.rgb = MID_GRAY


def add_title(slide, title, subtitle=None, dark=False):
    color = WHITE if dark else NAVY
    sub_color = RGBColor(214, 228, 245) if dark else MID_GRAY
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.8), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = TITLE_FONT
    run.font.bold = True
    run.font.size = Pt(34)
    run.font.color.rgb = color
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.82), Inches(1.35), Inches(10.5), Inches(0.45))
        stf = sub_box.text_frame
        sp = stf.paragraphs[0]
        sp.text = subtitle
        srun = sp.runs[0]
        srun.font.name = BODY_FONT
        srun.font.size = Pt(16)
        srun.font.color.rgb = sub_color


def add_card(slide, x, y, w, h, title, body, accent=TEAL, title_color=DARK):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = LINE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.18), Inches(w - 0.35), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    run = p.runs[0]
    run.font.name = TITLE_FONT
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = title_color

    bb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.56), Inches(w - 0.35), Inches(h - 0.7))
    btf = bb.text_frame
    btf.word_wrap = True
    for i, line in enumerate(body):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.text = line
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name = BODY_FONT
        run.font.size = Pt(14)
        run.font.color.rgb = MID_GRAY


def add_bullet_list(slide, x, y, w, h, items, size=18, color=DARK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        p.bullet = True
        run = p.runs[0]
        run.font.name = BODY_FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color


def add_kpi(slide, x, y, w, h, value, label, color):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = LINE

    v = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.18), Inches(w - 0.4), Inches(0.5))
    vp = v.text_frame.paragraphs[0]
    vp.text = value
    vr = vp.runs[0]
    vr.font.name = TITLE_FONT
    vr.font.bold = True
    vr.font.size = Pt(28)
    vr.font.color.rgb = color

    l = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.73), Inches(w - 0.4), Inches(0.35))
    lp = l.text_frame.paragraphs[0]
    lp.text = label
    lr = lp.runs[0]
    lr.font.name = BODY_FONT
    lr.font.size = Pt(13)
    lr.font.color.rgb = MID_GRAY


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = 14

    # 1 cover
    s = prs.slides.add_slide(blank)
    add_bg(s, NAVY)
    add_title(s, "知识管理培训", "10分钟快速上手：从“资料堆积”到“可复用资产”", dark=True)
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(2.2), Inches(4.5), Inches(0.55))
    pill.fill.solid()
    pill.fill.fore_color.rgb = TEAL
    pill.line.fill.background()
    tp = pill.text_frame.paragraphs[0]
    tp.text = "结构化 · 标准化 · 视觉化"
    tp.runs[0].font.name = BODY_FONT
    tp.runs[0].font.size = Pt(16)
    tp.runs[0].font.bold = True
    tp.runs[0].font.color.rgb = WHITE
    add_card(s, 0.85, 3.2, 3.9, 2.2, "培训收益", ["统一知识口径", "减少重复沟通", "加速新人上手"], accent=RGBColor(108, 190, 255))
    add_card(s, 4.95, 3.2, 3.9, 2.2, "适用对象", ["团队管理者", "业务骨干", "流程与运营岗位"], accent=RGBColor(134, 214, 192))
    add_card(s, 9.05, 3.2, 3.45, 2.2, "输出物", ["知识地图", "知识卡模板", "90天落地计划"], accent=RGBColor(255, 196, 118))
    add_footer(s, 1, total)

    # 2 agenda
    s = prs.slides.add_slide(blank)
    add_bg(s, LIGHT_BG)
    add_title(s, "培训目标与议程", "10分钟分为 4 段：认知、方法、落地、行动")
    stages = [("01 认知", "为什么做", NAVY), ("02 方法", "怎么做", TEAL), ("03 落地", "谁来做", GREEN), ("04 行动", "何时做", ORANGE)]
    x = 0.9
    for name, desc, c in stages:
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.2), Inches(2.8), Inches(2.0))
        rect.fill.solid()
        rect.fill.fore_color.rgb = WHITE
        rect.line.color.rgb = LINE
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.1), Inches(2.5), Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = c
        circle.line.fill.background()
        t = s.shapes.add_textbox(Inches(x + 0.25), Inches(3.25), Inches(2.3), Inches(0.8))
        tf = t.text_frame
        p1 = tf.paragraphs[0]
        p1.text = name
        p1.runs[0].font.name = BODY_FONT
        p1.runs[0].font.bold = True
        p1.runs[0].font.size = Pt(18)
        p1.runs[0].font.color.rgb = DARK
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.runs[0].font.name = BODY_FONT
        p2.runs[0].font.size = Pt(14)
        p2.runs[0].font.color.rgb = MID_GRAY
        x += 3.1
    add_footer(s, 2, total)

    # 3 why km
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    add_title(s, "为什么要做知识管理", "问题不是“没有知识”，而是“知识不可被快速找到和复用”")
    add_card(s, 0.8, 1.9, 6.1, 4.7, "典型痛点", ["同样问题反复问", "文档命名混乱，搜索困难", "离职/轮岗导致经验断层", "新人上手周期过长"], accent=RED)
    add_card(s, 6.95, 1.9, 5.55, 4.7, "预期收益", ["效率：检索时间下降 30%-50%", "质量：关键流程标准化", "组织：经验沉淀为可复用资产", "协同：跨团队沟通成本下降"], accent=GREEN)
    stat = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.25), Inches(5.7), Inches(5.0), Inches(0.8))
    stat.fill.solid()
    stat.fill.fore_color.rgb = MINT
    stat.line.fill.background()
    st = stat.text_frame.paragraphs[0]
    st.text = "核心原则：让正确的人，在正确时间，拿到正确知识"
    st.runs[0].font.name = BODY_FONT
    st.runs[0].font.size = Pt(14)
    st.runs[0].font.bold = True
    st.runs[0].font.color.rgb = TEAL
    add_footer(s, 3, total)

    # 4 model
    s = prs.slides.add_slide(blank)
    add_bg(s, LIGHT_BG)
    add_title(s, "知识管理全景模型", "形成闭环：采集 → 结构化 → 分发 → 应用 → 复盘")
    steps = [("采集", NAVY), ("结构化", TEAL), ("分发", GREEN), ("应用", ORANGE), ("复盘", RED)]
    x = 0.95
    for i, (name, c) in enumerate(steps):
        node = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(3.0), Inches(1.8), Inches(1.8))
        node.fill.solid()
        node.fill.fore_color.rgb = c
        node.line.fill.background()
        tp = node.text_frame.paragraphs[0]
        tp.text = name
        tp.alignment = PP_ALIGN.CENTER
        tp.runs[0].font.name = BODY_FONT
        tp.runs[0].font.size = Pt(16)
        tp.runs[0].font.bold = True
        tp.runs[0].font.color.rgb = WHITE
        if i < len(steps) - 1:
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 1.75), Inches(3.62), Inches(0.65), Inches(0.55))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(178, 190, 204)
            arrow.line.fill.background()
        x += 2.45
    add_bullet_list(s, 1.0, 5.35, 11.4, 1.7, ["每个环节都要有“责任人 + 产出标准 + 更新频率”", "闭环越短，知识越新鲜，复用价值越高"], size=16)
    add_footer(s, 4, total)

    # 5 pyramid
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    add_title(s, "方法一：结构化沉淀", "把“零散信息”加工成“可复用资产”")
    levels = [
        ("知识资产（模板/案例库）", TEAL, 8.3, 0.95),
        ("方法（步骤、准则、决策树）", RGBColor(53, 154, 154), 6.8, 2.1),
        ("知识点（定义、规则、结论）", RGBColor(86, 178, 178), 5.3, 3.25),
        ("信息（记录、聊天、原始数据）", RGBColor(121, 202, 202), 3.8, 4.4),
    ]
    for text, c, w, y in levels:
        x = (13.333 - w) / 2
        shape = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x), Inches(y), Inches(w), Inches(1.05))
        shape.fill.solid()
        shape.fill.fore_color.rgb = c
        shape.line.fill.background()
        p = shape.text_frame.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.name = BODY_FONT
        p.runs[0].font.size = Pt(14)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = WHITE
    add_footer(s, 5, total)

    # 6 knowledge card
    s = prs.slides.add_slide(blank)
    add_bg(s, LIGHT_BG)
    add_title(s, "方法二：标准化模板", "建议统一使用“一页知识卡”模板，降低撰写和阅读成本")
    outer = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.9))
    outer.fill.solid()
    outer.fill.fore_color.rgb = WHITE
    outer.line.color.rgb = LINE
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.8), Inches(11.3), Inches(0.9))
    top.fill.solid()
    top.fill.fore_color.rgb = NAVY
    top.line.fill.background()
    tt = top.text_frame.paragraphs[0]
    tt.text = "知识卡示例：客户投诉处理 SOP"
    tt.runs[0].font.name = BODY_FONT
    tt.runs[0].font.size = Pt(18)
    tt.runs[0].font.bold = True
    tt.runs[0].font.color.rgb = WHITE
    add_card(s, 1.4, 2.95, 3.2, 3.3, "场景", ["何时使用", "触发条件"], accent=TEAL)
    add_card(s, 4.95, 2.95, 3.2, 3.3, "步骤", ["1) 定级", "2) 响应", "3) 闭环"], accent=GREEN)
    add_card(s, 8.5, 2.95, 3.2, 3.3, "注意项", ["升级条件", "常见错误", "复盘清单"], accent=ORANGE)
    add_footer(s, 6, total)

    # 7 visualize
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    add_title(s, "方法三：视觉化表达", "同一知识可用三种图形表达：流程图、矩阵图、时间线")
    add_card(s, 0.9, 2.0, 4.05, 4.9, "流程图", ["适合：步骤明确", "价值：降低理解偏差", "示例：审批流程"], accent=NAVY)
    add_card(s, 4.98, 2.0, 4.05, 4.9, "矩阵图", ["适合：分类与优先级", "价值：辅助决策", "示例：重要/紧急四象限"], accent=TEAL)
    add_card(s, 9.06, 2.0, 3.35, 4.9, "时间线", ["适合：阶段计划", "价值：统一节奏", "示例：30-60-90天"], accent=GREEN)
    add_footer(s, 7, total)

    # 8 tool map
    s = prs.slides.add_slide(blank)
    add_bg(s, LIGHT_BG)
    add_title(s, "工具与载体地图", "不追求工具多，追求“入口统一 + 版本唯一 + 权责明确”")
    frame = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(2.0), Inches(10.9), Inches(4.8))
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = LINE
    vline = s.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, Inches(6.65), Inches(2.0), Inches(0), Inches(4.8))
    vline.line.color.rgb = LINE
    hline = s.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, Inches(1.2), Inches(4.4), Inches(10.9), Inches(0))
    hline.line.color.rgb = LINE
    quads = [
        (1.45, 2.2, "文档库", "制度、模板、归档"),
        (6.85, 2.2, "Wiki", "知识主题页、FAQ"),
        (1.45, 4.6, "问答库", "高频问题闭环"),
        (6.85, 4.6, "专家网络", "谁知道、找谁问"),
    ]
    for x, y, t, d in quads:
        box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(4.9), Inches(1.6))
        tf = box.text_frame
        p1 = tf.paragraphs[0]
        p1.text = t
        p1.runs[0].font.name = BODY_FONT
        p1.runs[0].font.bold = True
        p1.runs[0].font.size = Pt(22)
        p1.runs[0].font.color.rgb = NAVY
        p2 = tf.add_paragraph()
        p2.text = d
        p2.runs[0].font.name = BODY_FONT
        p2.runs[0].font.size = Pt(14)
        p2.runs[0].font.color.rgb = MID_GRAY
    add_footer(s, 8, total)

    # 9 process
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    add_title(s, "团队落地流程（7步）", "每一步都要可交付、可检查、可迭代")
    steps = ["盘点", "分层", "建模", "模板", "发布", "运营", "复盘"]
    colors = [NAVY, RGBColor(27, 78, 126), TEAL, RGBColor(8, 146, 146), GREEN, ORANGE, RED]
    x = 0.75
    for i, name in enumerate(steps):
        w = 1.65
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.0), Inches(w), Inches(1.0))
        rect.fill.solid()
        rect.fill.fore_color.rgb = colors[i]
        rect.line.fill.background()
        p = rect.text_frame.paragraphs[0]
        p.text = f"{i+1}. {name}"
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.name = BODY_FONT
        p.runs[0].font.size = Pt(14)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = WHITE
        if i < len(steps) - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x + 1.62), Inches(3.15), Inches(0.34), Inches(0.7))
            arr.fill.solid()
            arr.fill.fore_color.rgb = RGBColor(184, 195, 208)
            arr.line.fill.background()
        x += 1.82
    add_bullet_list(s, 1.0, 4.55, 11.3, 2.0, ["建议先选一个业务场景做试点，2周验证后再扩面", "流程可复制，模板需按团队场景定制"], size=16)
    add_footer(s, 9, total)

    # 10 roles
    s = prs.slides.add_slide(blank)
    add_bg(s, LIGHT_BG)
    add_title(s, "角色分工（谁来做）", "用 RACI 思路明确责任，避免“都管=没人管”")
    headers = ["角色", "R 负责", "A 负责到底", "C 协作", "I 知会"]
    rows = [
        ["业务专家", "提供知识源", "", "评审标准", "同步变更"],
        ["知识管理员", "整理发布", "版本管理", "推动更新", "广播通知"],
        ["团队负责人", "", "目标与资源", "复盘节奏", "绩效关联"],
    ]
    x0, y0, total_w, row_h = 1.0, 2.1, 11.3, 1.05
    col_w = [2.0, 2.2, 2.4, 2.3, 2.4]
    x = x0
    for i, h in enumerate(headers):
        c = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y0), Inches(col_w[i]), Inches(0.85))
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
        c.line.color.rgb = WHITE
        p = c.text_frame.paragraphs[0]
        p.text = h
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.name = BODY_FONT
        p.runs[0].font.bold = True
        p.runs[0].font.size = Pt(13)
        p.runs[0].font.color.rgb = WHITE
        x += col_w[i]
    for r, row in enumerate(rows):
        x = x0
        for i, txt in enumerate(row):
            cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y0 + 0.85 + r * row_h), Inches(col_w[i]), Inches(row_h))
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 == 0 else RGBColor(248, 251, 253)
            cell.line.color.rgb = LINE
            p = cell.text_frame.paragraphs[0]
            p.text = txt
            p.alignment = PP_ALIGN.CENTER
            if p.runs:
                p.runs[0].font.name = BODY_FONT
                p.runs[0].font.size = Pt(12)
                p.runs[0].font.color.rgb = MID_GRAY
            x += col_w[i]
    add_footer(s, 10, total)

    # 11 KPI
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    add_title(s, "衡量指标（怎么判断有效）", "同时关注“使用量、质量、时效、业务结果”")
    add_kpi(s, 0.9, 1.95, 2.9, 1.3, "75%", "知识检索命中率", TEAL)
    add_kpi(s, 3.95, 1.95, 2.9, 1.3, "48h", "知识更新时效", ORANGE)
    add_kpi(s, 7.0, 1.95, 2.9, 1.3, "62%", "模板使用覆盖率", NAVY)
    add_kpi(s, 10.05, 1.95, 2.4, 1.3, "-30%", "重复问答量", GREEN)
    bars = [("试点前", 2.3, RGBColor(183, 196, 212)), ("第30天", 3.6, RGBColor(123, 175, 214)), ("第60天", 4.5, RGBColor(73, 158, 189)), ("第90天", 5.2, TEAL)]
    y = 4.0
    for label, w, c in bars:
        lab = s.shapes.add_textbox(Inches(1.1), Inches(y + 0.07), Inches(1.3), Inches(0.35))
        p = lab.text_frame.paragraphs[0]
        p.text = label
        p.runs[0].font.name = BODY_FONT
        p.runs[0].font.size = Pt(12)
        p.runs[0].font.color.rgb = MID_GRAY
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.6), Inches(y), Inches(w), Inches(0.45))
        bar.fill.solid()
        bar.fill.fore_color.rgb = c
        bar.line.fill.background()
        y += 0.75
    add_footer(s, 11, total)

    # 12 mistakes
    s = prs.slides.add_slide(blank)
    add_bg(s, LIGHT_BG)
    add_title(s, "常见误区与纠偏", "先避免踩坑，再追求复杂功能")
    add_card(s, 0.9, 2.0, 5.95, 4.9, "常见误区", ["重建设，轻运营", "重工具，轻标准", "重采集，轻复盘", "重数量，轻质量"], accent=RED)
    add_card(s, 6.5, 2.0, 5.95, 4.9, "纠偏动作", ["以业务场景为起点", "先模板后平台", "设定固定复盘节奏", "把指标绑定到团队目标"], accent=GREEN)
    add_footer(s, 12, total)

    # 13 30-60-90
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    add_title(s, "30-60-90 天行动计划", "用可执行计划替代“知道但不做”")
    plan = [("0-30天", NAVY, ["盘点高频问题", "确定模板标准", "完成试点范围"]),
            ("31-60天", TEAL, ["上线知识地图", "推行每周更新", "建立评审机制"]),
            ("61-90天", GREEN, ["扩大到跨团队", "看板追踪指标", "形成季度复盘"])]
    x = 0.95
    for title, c, items in plan:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.2), Inches(3.95), Inches(4.6))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LINE
        head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.2), Inches(3.95), Inches(0.8))
        head.fill.solid()
        head.fill.fore_color.rgb = c
        head.line.fill.background()
        p = head.text_frame.paragraphs[0]
        p.text = title
        p.runs[0].font.name = BODY_FONT
        p.runs[0].font.bold = True
        p.runs[0].font.size = Pt(18)
        p.runs[0].font.color.rgb = WHITE
        add_bullet_list(s, x + 0.2, 3.25, 3.5, 3.2, items, size=14)
        if x < 8.0:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 3.98), Inches(4.05), Inches(0.55), Inches(0.55))
            ar.fill.solid()
            ar.fill.fore_color.rgb = RGBColor(190, 201, 214)
            ar.line.fill.background()
        x += 4.2
    add_footer(s, 13, total)

    # 14 ending
    s = prs.slides.add_slide(blank)
    add_bg(s, NAVY)
    add_title(s, "总结：知识管理是“业务能力工程”", "小步快跑，先做起来，再持续优化", dark=True)
    add_bullet_list(s, 1.0, 2.2, 10.8, 2.7, ["先选一个场景试点：两周出结果", "统一一页知识卡：降低沟通成本", "建立周更新 + 月复盘：保持知识新鲜度"], size=22, color=WHITE)
    end = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.3), Inches(11.3), Inches(1.2))
    end.fill.solid()
    end.fill.fore_color.rgb = TEAL
    end.line.fill.background()
    p = end.text_frame.paragraphs[0]
    p.text = "行动口号：把经验变成资产，把资产变成效率"
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.name = TITLE_FONT
    p.runs[0].font.bold = True
    p.runs[0].font.size = Pt(24)
    p.runs[0].font.color.rgb = WHITE
    add_footer(s, 14, total)

    output = "知识管理培训_10分钟_14页.pptx"
    prs.save(output)
    print(output)


if __name__ == "__main__":
    build()
