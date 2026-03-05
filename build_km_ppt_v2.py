from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


FONT = "PingFang SC"

BG_LIGHT = RGBColor(248, 250, 252)
BG_DARK = RGBColor(15, 23, 42)
PRIMARY = RGBColor(14, 116, 144)
PRIMARY_2 = RGBColor(13, 148, 136)
ACCENT = RGBColor(245, 158, 11)
TEXT = RGBColor(15, 23, 42)
MUTED = RGBColor(100, 116, 139)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(226, 232, 240)
OK = RGBColor(22, 163, 74)
WARN = RGBColor(234, 88, 12)
BAD = RGBColor(220, 38, 38)

W = Inches(13.333)
H = Inches(7.5)


def bg(slide, color=BG_LIGHT):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()


def deco_light(slide):
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.9), Inches(0.0), Inches(2.3), Inches(2.3))
    c1.fill.solid()
    c1.fill.fore_color.rgb = RGBColor(224, 242, 254)
    c1.line.fill.background()
    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.0), Inches(5.3), Inches(1.9), Inches(1.9))
    c2.fill.solid()
    c2.fill.fore_color.rgb = RGBColor(204, 251, 241)
    c2.line.fill.background()


def title(slide, t, s=None, dark=False):
    tc = WHITE if dark else TEXT
    sc = RGBColor(203, 213, 225) if dark else MUTED
    tb = slide.shapes.add_textbox(Inches(0.85), Inches(0.45), Inches(11.5), Inches(0.9))
    p = tb.text_frame.paragraphs[0]
    p.text = t
    p.runs[0].font.name = FONT
    p.runs[0].font.bold = True
    p.runs[0].font.size = Pt(36)
    p.runs[0].font.color.rgb = tc
    if s:
        sb = slide.shapes.add_textbox(Inches(0.87), Inches(1.35), Inches(11.0), Inches(0.45))
        sp = sb.text_frame.paragraphs[0]
        sp.text = s
        sp.runs[0].font.name = FONT
        sp.runs[0].font.size = Pt(16)
        sp.runs[0].font.color.rgb = sc


def footer(slide, page, total, dark=False):
    c = RGBColor(148, 163, 184) if dark else RGBColor(148, 163, 184)
    fb = slide.shapes.add_textbox(Inches(11.9), Inches(7.05), Inches(1.1), Inches(0.25))
    p = fb.text_frame.paragraphs[0]
    p.text = f"{page:02d}/{total:02d}"
    p.alignment = PP_ALIGN.RIGHT
    p.runs[0].font.name = FONT
    p.runs[0].font.size = Pt(11)
    p.runs[0].font.color.rgb = c


def chip(slide, x, y, text, fill=PRIMARY, color=WHITE):
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.6), Inches(0.46))
    r.fill.solid()
    r.fill.fore_color.rgb = fill
    r.line.fill.background()
    p = r.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.name = FONT
    p.runs[0].font.bold = True
    p.runs[0].font.size = Pt(12)
    p.runs[0].font.color.rgb = color


def card(slide, x, y, w, h, head, lines, head_color=TEXT, bar=PRIMARY):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid()
    c.fill.fore_color.rgb = WHITE
    c.line.color.rgb = LINE
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.09), Inches(h))
    b.fill.solid()
    b.fill.fore_color.rgb = bar
    b.line.fill.background()
    hb = slide.shapes.add_textbox(Inches(x + 0.24), Inches(y + 0.16), Inches(w - 0.35), Inches(0.35))
    hp = hb.text_frame.paragraphs[0]
    hp.text = head
    hp.runs[0].font.name = FONT
    hp.runs[0].font.bold = True
    hp.runs[0].font.size = Pt(18)
    hp.runs[0].font.color.rgb = head_color
    bb = slide.shapes.add_textbox(Inches(x + 0.24), Inches(y + 0.62), Inches(w - 0.4), Inches(h - 0.75))
    tf = bb.text_frame
    tf.word_wrap = True
    for i, t in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t
        p.bullet = True
        p.space_after = Pt(7)
        p.runs[0].font.name = FONT
        p.runs[0].font.size = Pt(14)
        p.runs[0].font.color.rgb = MUTED


def section(slide, name, subtitle):
    bg(slide, BG_DARK)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.4), Inches(13.333), Inches(2.2))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(30, 41, 59)
    band.line.fill.background()
    t = slide.shapes.add_textbox(Inches(1.0), Inches(2.95), Inches(11.3), Inches(0.8))
    p = t.text_frame.paragraphs[0]
    p.text = name
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.name = FONT
    p.runs[0].font.bold = True
    p.runs[0].font.size = Pt(40)
    p.runs[0].font.color.rgb = WHITE
    s = slide.shapes.add_textbox(Inches(1.0), Inches(3.8), Inches(11.3), Inches(0.45))
    sp = s.text_frame.paragraphs[0]
    sp.text = subtitle
    sp.alignment = PP_ALIGN.CENTER
    sp.runs[0].font.name = FONT
    sp.runs[0].font.size = Pt(16)
    sp.runs[0].font.color.rgb = RGBColor(148, 163, 184)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]
    total = 12

    # 1 cover
    s = prs.slides.add_slide(blank)
    bg(s, BG_DARK)
    d1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.8), Inches(0.0), Inches(3.4), Inches(3.4))
    d1.fill.solid()
    d1.fill.fore_color.rgb = RGBColor(15, 118, 110)
    d1.line.fill.background()
    d2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.0), Inches(5.1), Inches(2.4), Inches(2.4))
    d2.fill.solid()
    d2.fill.fore_color.rgb = RGBColor(3, 105, 161)
    d2.line.fill.background()
    title(s, "知识管理培训", "10分钟 | 结构化沉淀 · 可视化表达 · 业务化落地", dark=True)
    chip(s, 0.9, 2.2, "简单", PRIMARY_2)
    chip(s, 3.7, 2.2, "结构化", PRIMARY)
    chip(s, 6.5, 2.2, "视觉化", ACCENT, TEXT)
    big = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(3.1), Inches(11.6), Inches(2.8))
    big.fill.solid()
    big.fill.fore_color.rgb = RGBColor(30, 41, 59)
    big.line.fill.background()
    tb = s.shapes.add_textbox(Inches(1.25), Inches(3.55), Inches(10.9), Inches(2.0))
    tf = tb.text_frame
    for i, tx in enumerate(
        [
            "把分散经验，变成可检索、可复用、可传承的组织资产",
            "从“找不到”到“找得到”，从“看不懂”到“用得上”",
            "从“个人经验”到“团队能力”",
        ]
    ):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = tx
        p.space_after = Pt(8)
        p.runs[0].font.name = FONT
        p.runs[0].font.size = Pt(20 if i == 0 else 17)
        p.runs[0].font.bold = i == 0
        p.runs[0].font.color.rgb = WHITE if i == 0 else RGBColor(203, 213, 225)
    footer(s, 1, total, True)

    # 2 agenda
    s = prs.slides.add_slide(blank)
    bg(s)
    deco_light(s)
    title(s, "培训结构", "认知 -> 方法 -> 落地 -> 行动")
    labels = [("01 认知", "为什么做"), ("02 方法", "怎么做"), ("03 落地", "谁来做"), ("04 行动", "何时做")]
    x = 1.0
    for i, (a, b) in enumerate(labels):
        c = [PRIMARY, PRIMARY_2, OK, ACCENT][i]
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.9), Inches(2.45), Inches(0.55), Inches(0.55))
        dot.fill.solid()
        dot.fill.fore_color.rgb = c
        dot.line.fill.background()
        if i < 3:
            l = s.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, Inches(x + 1.45), Inches(2.73), Inches(2.0), Inches(0))
            l.line.color.rgb = RGBColor(186, 230, 253)
        tx = s.shapes.add_textbox(Inches(x), Inches(3.2), Inches(2.5), Inches(1.2))
        tf = tx.text_frame
        p1 = tf.paragraphs[0]
        p1.text = a
        p1.alignment = PP_ALIGN.CENTER
        p1.runs[0].font.name = FONT
        p1.runs[0].font.bold = True
        p1.runs[0].font.size = Pt(20)
        p1.runs[0].font.color.rgb = TEXT
        p2 = tf.add_paragraph()
        p2.text = b
        p2.alignment = PP_ALIGN.CENTER
        p2.runs[0].font.name = FONT
        p2.runs[0].font.size = Pt(14)
        p2.runs[0].font.color.rgb = MUTED
        x += 3.05
    footer(s, 2, total)

    # 3 section
    s = prs.slides.add_slide(blank)
    section(s, "01 认知", "从痛点出发，建立知识管理共识")
    footer(s, 3, total, True)

    # 4 why
    s = prs.slides.add_slide(blank)
    bg(s)
    deco_light(s)
    title(s, "为什么要做知识管理", "真正的问题是“知识无法被稳定复用”")
    card(s, 0.9, 1.95, 5.9, 4.9, "现状痛点", ["问题反复出现，重复沟通多", "知识散落在聊天/网盘/个人电脑", "新人培养依赖“师傅带徒弟”", "关键经验随人员流动而丢失"], bar=BAD)
    card(s, 6.55, 1.95, 5.9, 4.9, "业务价值", ["检索效率提升，减少无效等待", "流程标准化，质量波动更小", "经验可传承，组织韧性更强", "跨团队协作更顺畅"], bar=OK)
    footer(s, 4, total)

    # 5 model
    s = prs.slides.add_slide(blank)
    bg(s, WHITE)
    title(s, "知识管理闭环", "采集 -> 结构化 -> 分发 -> 应用 -> 复盘")
    steps = [("采集", PRIMARY), ("结构化", PRIMARY_2), ("分发", OK), ("应用", ACCENT), ("复盘", WARN)]
    x = 1.2
    for i, (name, c) in enumerate(steps):
        node = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.6), Inches(1.75), Inches(1.75))
        node.fill.solid()
        node.fill.fore_color.rgb = c
        node.line.fill.background()
        p = node.text_frame.paragraphs[0]
        p.text = name
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.name = FONT
        p.runs[0].font.bold = True
        p.runs[0].font.size = Pt(16)
        p.runs[0].font.color.rgb = WHITE if c != ACCENT else TEXT
        if i < 4:
            arr = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x + 1.72), Inches(3.1), Inches(0.45), Inches(0.7))
            arr.fill.solid()
            arr.fill.fore_color.rgb = RGBColor(203, 213, 225)
            arr.line.fill.background()
        x += 2.35
    note = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.1), Inches(11.3), Inches(1.2))
    note.fill.solid()
    note.fill.fore_color.rgb = RGBColor(241, 245, 249)
    note.line.fill.background()
    np = note.text_frame.paragraphs[0]
    np.text = "闭环要素：每个环节都设置责任人、输出标准、更新时间，避免“建而不用”。"
    np.runs[0].font.name = FONT
    np.runs[0].font.size = Pt(16)
    np.runs[0].font.color.rgb = TEXT
    footer(s, 5, total)

    # 6 method
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "三种方法，快速落地", "先统一表达方式，再推动规模化沉淀")
    card(s, 0.9, 2.0, 4.05, 4.75, "结构化", ["按主题/场景/流程分层", "把信息加工为知识卡", "形成可检索目录"], bar=PRIMARY)
    card(s, 4.98, 2.0, 4.05, 4.75, "标准化", ["统一模板与命名规则", "明确版本和负责人", "设定更新频率"], bar=PRIMARY_2)
    card(s, 9.06, 2.0, 3.35, 4.75, "视觉化", ["流程图", "矩阵图", "时间线"], bar=ACCENT)
    footer(s, 6, total)

    # 7 tool
    s = prs.slides.add_slide(blank)
    bg(s, WHITE)
    title(s, "工具不是重点，治理才是重点", "入口统一、版本唯一、职责清晰")
    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(1.95), Inches(11.1), Inches(4.9))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(248, 250, 252)
    panel.line.color.rgb = LINE
    labels = [
        (1.5, 2.45, "文档库", "制度/模板/标准"),
        (6.95, 2.45, "Wiki", "主题页/FAQ"),
        (1.5, 4.6, "问答库", "高频问题闭环"),
        (6.95, 4.6, "专家网络", "谁知道找谁"),
    ]
    for x, y, t1, t2 in labels:
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.75), Inches(1.55))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = LINE
        tp = box.text_frame
        p1 = tp.paragraphs[0]
        p1.text = t1
        p1.runs[0].font.name = FONT
        p1.runs[0].font.bold = True
        p1.runs[0].font.size = Pt(22)
        p1.runs[0].font.color.rgb = TEXT
        p2 = tp.add_paragraph()
        p2.text = t2
        p2.runs[0].font.name = FONT
        p2.runs[0].font.size = Pt(14)
        p2.runs[0].font.color.rgb = MUTED
    footer(s, 7, total)

    # 8 section
    s = prs.slides.add_slide(blank)
    section(s, "02 落地", "角色、指标、节奏三件事")
    footer(s, 8, total, True)

    # 9 role
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "角色分工（RACI）", "都参与，但责任必须唯一")
    card(s, 0.9, 2.05, 3.9, 4.8, "业务专家 R", ["提供一线经验", "评审知识准确性", "反馈应用问题"], bar=PRIMARY)
    card(s, 4.95, 2.05, 3.9, 4.8, "知识管理员 A", ["组织沉淀与发布", "维护版本和目录", "推动更新节奏"], bar=PRIMARY_2)
    card(s, 9.0, 2.05, 3.45, 4.8, "团队负责人 C/I", ["分配资源与目标", "跟踪效果指标", "在例会中持续复盘"], bar=ACCENT)
    footer(s, 9, total)

    # 10 kpi
    s = prs.slides.add_slide(blank)
    bg(s, WHITE)
    title(s, "指标看板（建议）", "看“是否被用”，再看“是否产生业务价值”")
    kpis = [("75%", "检索命中率", PRIMARY), ("48h", "知识更新时效", ACCENT), ("62%", "模板覆盖率", PRIMARY_2), ("-30%", "重复问答量", OK)]
    x_positions = [0.9, 3.8, 6.7, 9.6]
    for i, (v, l, c) in enumerate(kpis):
        x = x_positions[i]
        k = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.95), Inches(2.75), Inches(1.35))
        k.fill.solid()
        k.fill.fore_color.rgb = RGBColor(248, 250, 252)
        k.line.color.rgb = LINE
        vp = s.shapes.add_textbox(Inches(x + 0.2), Inches(2.12), Inches(2.35), Inches(0.5)).text_frame.paragraphs[0]
        vp.text = v
        vp.runs[0].font.name = FONT
        vp.runs[0].font.bold = True
        vp.runs[0].font.size = Pt(30)
        vp.runs[0].font.color.rgb = c
        lp = s.shapes.add_textbox(Inches(x + 0.2), Inches(2.72), Inches(2.35), Inches(0.3)).text_frame.paragraphs[0]
        lp.text = l
        lp.runs[0].font.name = FONT
        lp.runs[0].font.size = Pt(13)
        lp.runs[0].font.color.rgb = MUTED
    bars = [("试点前", 2.2, RGBColor(203, 213, 225)), ("30天", 3.3, RGBColor(125, 211, 252)), ("60天", 4.2, RGBColor(45, 212, 191)), ("90天", 5.0, PRIMARY)]
    y = 4.15
    for name, w, c in bars:
        lb = s.shapes.add_textbox(Inches(1.1), Inches(y + 0.08), Inches(1.2), Inches(0.3))
        p = lb.text_frame.paragraphs[0]
        p.text = name
        p.runs[0].font.name = FONT
        p.runs[0].font.size = Pt(12)
        p.runs[0].font.color.rgb = MUTED
        b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.35), Inches(y), Inches(w), Inches(0.44))
        b.fill.solid()
        b.fill.fore_color.rgb = c
        b.line.fill.background()
        y += 0.7
    footer(s, 10, total)

    # 11 roadmap
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "30-60-90 天行动路线", "先试点，再扩面，最后固化机制")
    stages = [
        ("0-30天", PRIMARY, ["完成知识盘点", "确定模板与命名", "跑通1个试点流程"]),
        ("31-60天", PRIMARY_2, ["上线知识地图", "建立每周更新", "设置质量评审"]),
        ("61-90天", OK, ["跨团队推广", "看板跟踪指标", "形成月度复盘机制"]),
    ]
    x = 0.95
    for i, (n, c, items) in enumerate(stages):
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.1), Inches(3.95), Inches(4.85))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = LINE
        head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.1), Inches(3.95), Inches(0.8))
        head.fill.solid()
        head.fill.fore_color.rgb = c
        head.line.fill.background()
        p = head.text_frame.paragraphs[0]
        p.text = n
        p.runs[0].font.name = FONT
        p.runs[0].font.bold = True
        p.runs[0].font.size = Pt(19)
        p.runs[0].font.color.rgb = WHITE
        t = s.shapes.add_textbox(Inches(x + 0.2), Inches(3.12), Inches(3.5), Inches(3.6)).text_frame
        for j, it in enumerate(items):
            pp = t.paragraphs[0] if j == 0 else t.add_paragraph()
            pp.text = it
            pp.bullet = True
            pp.space_after = Pt(8)
            pp.runs[0].font.name = FONT
            pp.runs[0].font.size = Pt(14)
            pp.runs[0].font.color.rgb = MUTED
        if i < 2:
            a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 4.0), Inches(4.2), Inches(0.5), Inches(0.5))
            a.fill.solid()
            a.fill.fore_color.rgb = RGBColor(203, 213, 225)
            a.line.fill.background()
        x += 4.2
    footer(s, 11, total)

    # 12 ending
    s = prs.slides.add_slide(blank)
    bg(s, BG_DARK)
    t1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.2), Inches(0.0), Inches(3.0), Inches(3.0))
    t1.fill.solid()
    t1.fill.fore_color.rgb = RGBColor(13, 148, 136)
    t1.line.fill.background()
    title(s, "把经验变成资产，把资产变成效率", "知识管理不是文档工程，而是业务能力工程", dark=True)
    card_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.7))
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = RGBColor(30, 41, 59)
    card_bg.line.fill.background()
    lines = [
        "从一个高频业务场景开始试点（2周见结果）",
        "统一一页知识卡模板（先标准，再规模）",
        "建立周更新 + 月复盘机制（保证长期有效）",
    ]
    tx = s.shapes.add_textbox(Inches(1.45), Inches(2.75), Inches(10.6), Inches(2.7)).text_frame
    for i, l in enumerate(lines):
        p = tx.paragraphs[0] if i == 0 else tx.add_paragraph()
        p.text = l
        p.bullet = True
        p.space_after = Pt(10)
        p.runs[0].font.name = FONT
        p.runs[0].font.size = Pt(21)
        p.runs[0].font.color.rgb = WHITE
    footer(s, 12, total, True)

    out = "知识管理培训_美化版_12页.pptx"
    prs.save(out)
    print(out)


if __name__ == "__main__":
    build()
